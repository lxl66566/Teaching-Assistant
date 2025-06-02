from pathlib import Path

import pymupdf4llm
from docx import Document
from loguru import logger
from pptx import Presentation

from ..utils.pdf import is_text_pdf
from .document_ocr import ocr_pdf_pages


def flatten(xss):
    return [x for xs in xss for x in xs]


def ocr_pdf(pdf_path: str | Path) -> str:
    """
    对 PDF 文件进行 OCR 处理，并返回 OCR 结果。

    Args:
        pdf_path: PDF 文件的路径。

    Returns:
        提取到的文本字符串（Markdown 格式），或者一个描述错误的字符串（如果处理失败）。
    """
    ocr_result = ocr_pdf_pages(str(pdf_path))
    return "\n".join(flatten(ocr_result.values()))


def process_pdf_and_get_text(pdf_path: str | Path) -> str:
    """
    处理一个 PDF 文件，并提取其文本内容。
    如果 PDF 是文字 PDF，则直接提取文本；否则进行 OCR。

    Args:
        pdf_path: PDF 文件的路径。

    Returns:
        提取到的文本字符串（Markdown 格式），或者一个描述错误的字符串（如果处理失败）。
    """
    if isinstance(pdf_path, str):
        pdf_path = Path(pdf_path)

    if not pdf_path.exists():
        raise FileNotFoundError(f"文件未找到 -> {pdf_path}")

    try:
        # 步骤 1: 判断是否为文字 PDF
        logger.info(f"正在判断文件 '{pdf_path.name}' 是否为文字 PDF...")
        is_text = is_text_pdf(pdf_path)  # is_text_pdf 函数内部会处理其自身的错误和打印

        if is_text:
            logger.info(f"文件 '{pdf_path.name}' 被判断为文字 PDF，直接提取文本。")
            return pymupdf4llm.to_markdown(pdf_path)
        else:
            logger.info(f"文件 '{pdf_path.name}' 未被判断为文字 PDF，执行 OCR。")
            return ocr_pdf(pdf_path)

    except Exception as e:
        error_msg = f"错误: 处理 PDF 文件 '{pdf_path}' 时发生未捕获的错误: {e}"
        logger.error(error_msg)
        raise


def extract_text_from_docx(docx_file_path: str):
    """
    从DOCX文件中提取所有文本。
    """
    document = Document(docx_file_path)  # [11]
    full_text = []
    for para in document.paragraphs:
        full_text.append(para.text)

    return "\n".join(full_text)


def extract_text_from_pptx(pptx_file_path: str):
    """
    从PPTX文件中提取所有文本，包括幻灯片文本和备注。
    """
    prs = Presentation(pptx_file_path)  # [3]
    full_text = []
    for slide in prs.slides:
        # 提取幻灯片上的文本框内容
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:  # type: ignore # [13]
                text = shape.text_frame.text  # type: ignore
                if text.strip():
                    full_text.append(text)

        # 提取演示者备注 (如果有的话)
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide  # [3]
            if notes_slide.notes_text_frame:  # [6]
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    full_text.append(f"备注: {notes_text}")

    return "\n".join(full_text)


def process_file_to_text(file_path: str | Path, file_type: str):
    """
    处理文件，并提取其文本内容。
    如果文件是 DOCX、PPTX、ppt, doc 格式，则会调用对应的提取函数；如果文件是 txt, markdown 格式，则直接返回文件内容。
    如果文件是 pdf 则判断是否需要 ocr。
    """
    if isinstance(file_path, str):
        file_path = Path(file_path)

    if not file_path.exists():
        raise FileNotFoundError(f"文件未找到 -> {file_path}")

    try:
        # 步骤 1: 判断文件类型
        logger.info(f"正在判断文件 '{file_path.name}' 类型...")
        if file_type in ["docx", "doc"]:
            logger.info(f"文件 '{file_path.name}' 是 DOCX 格式，将调用对应的提取函数。")
            return extract_text_from_docx(str(file_path))
        elif file_type in ["pptx", "ppt"]:
            logger.info(f"文件 '{file_path.name}' 是 PPTX 格式，将调用对应的提取函数。")
            return extract_text_from_pptx(str(file_path))
        elif file_type in ["txt", "markdown"]:
            logger.info(f"文件 '{file_path.name}' 是纯文本格式，直接返回文件内容。")
            return file_path.read_text(encoding="utf-8")
        elif file_type == "pdf":
            logger.info(f"文件 '{file_path.name}' 是 PDF 格式，将进行 OCR。")
            return ocr_pdf(file_path)
        else:
            raise ValueError(f"文件 '{file_path.name}' 类型不支持。")
    except Exception as e:
        error_msg = f"错误: 处理文件 '{file_path}' 时发生错误: {e}"
        logger.error(error_msg)
        raise
